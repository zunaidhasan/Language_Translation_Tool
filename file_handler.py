from PIL import Image
import pytesseract
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_text_from_image(image_path):
    try:
        return pytesseract.image_to_string(Image.open(image_path))
    except Exception as e:
        return f"Error extracting text from image: {str(e)}"

def extract_text_from_pdf(pdf_path):
    try:
        reader = PdfReader(pdf_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text()
        return text
    except Exception as e:
        return f"Error extracting text from PDF: {str(e)}"

def extract_text_from_ppt(ppt_path):
    try:
        presentation = Presentation(ppt_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting text from PowerPoint: {str(e)}"
